import argparse
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

# --------------------
# CLI
# --------------------
parser = argparse.ArgumentParser()
parser.add_argument("--project", required=True)
parser.add_argument("--audience", required=True)
args = parser.parse_args()

# --------------------
# Paths
# --------------------
BASE = Path(f"projects/{args.project}")
MD = BASE / "deck/deck.md"
CHARTS = BASE / "financial/charts"
OUT = Path(f"outputs/{args.project}-{args.audience}.pptx")
OUT.parent.mkdir(exist_ok=True)

# --------------------
# Load branded template
# --------------------
prs = Presentation("templates/Brand_Template.pptx")

# --------------------
# Layout mapping (ADJUST TO YOUR TEMPLATE)
# --------------------
COVER = prs.slide_layouts[0]
SECTION = prs.slide_layouts[1]
CONTENT = prs.slide_layouts[3]
FINANCIAL = prs.slide_layouts[4]

# --------------------
# Helpers
# --------------------
def split_bullet_and_source(text: str):
    if "(Source:" in text:
        main, source = text.split("(Source:", 1)
        return main.strip(), source.replace(")", "").strip()
    return text, None

def add_footer(slide, text):
    footer = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(6.8),
        Inches(9),
        Inches(0.4)
    )
    footer.text_frame.text = text

# --------------------
# Parse markdown
# --------------------
content = MD.read_text(encoding="utf-8")
slides = content.split("# ")[1:]

first_slide = True

for slide in slides:
    lines = slide.strip().split("\n")
    title = lines[0].strip()
    bullets = [l.replace("- ", "").strip() for l in lines if l.startswith("-")]
    title_lower = title.lower()

    # --------------------
    # Cover
    # --------------------
    if first_slide:
        s = prs.slides.add_slide(COVER)
        s.shapes.title.text = title
        first_slide = False
        continue

    # --------------------
    # Financial slide
    # --------------------
    if "financial" in title_lower:
        s = prs.slides.add_slide(FINANCIAL)
        s.shapes.title.text = title

        chart_path = CHARTS / "revenue.png"
        if chart_path.exists():
            s.shapes.add_picture(
                str(chart_path),
                Inches(1),
                Inches(1.8),
                width=Inches(8)
            )

        add_footer(s, "Source: Financial_Model_Template.xlsx")
        continue

    # --------------------
    # Section slide
    # --------------------
    if not bullets:
        s = prs.slides.add_slide(SECTION)
        s.shapes.title.text = title
        continue

    # --------------------
    # Content slide
    # --------------------
    s = prs.slides.add_slide(CONTENT)
    s.shapes.title.text = title
    body = s.placeholders[1]
    body.text = ""

    sources = set()

    for i, bullet in enumerate(bullets):
        text, source = split_bullet_and_source(bullet)
        if source:
            sources.add(source)

        body.text += ("" if i == 0 else "\n") + text

    if sources:
        add_footer(s, f"Source: {', '.join(sorted(sources))}")

# --------------------
# Save
# --------------------
prs.save(OUT)
print(f"Branded PPT generated: {OUT}")